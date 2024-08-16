from dotenv import load_dotenv
import os
import sqlite3
import streamlit as st
from userDB import UserAuth
from datetime import datetime
import requests
from langchain.embeddings import HuggingFaceEmbeddings as hfe
import streamlit as st
from langchain_community.llms import Ollama
from langchain.text_splitter import RecursiveCharacterTextSplitter as rcts
from langchain_community.embeddings.fastembed import FastEmbedEmbeddings as fee
from langchain_community.document_loaders import PDFPlumberLoader as pdfPL
from langchain_community.vectorstores import Chroma
from langchain.chains.combine_documents import create_stuff_documents_chain as csdc
from langchain.chains import create_retrieval_chain as crc
from langchain.prompts import PromptTemplate as pt
from docx import Document as DocxDocument
from langchain.docstore.document import Document
import pandas as pd
from pptx import Presentation

# Set up Streamlit page
st.set_page_config(page_title="RAG Streamlit Application", layout="wide")

# Initialize global variables
llm = Ollama(model="llama3")

# Replace the embedding initialization
embedding = hfe(model_name="bert-base-uncased")

text_splitter = rcts(
    chunk_size=2048,
    chunk_overlap=200,
    length_function=len,
    is_separator_regex='\n'
)

# Ensure the database directory exists or create it
dbDirectory = r".\CorporateGPT2\db"
if not os.path.exists(dbDirectory):
    os.makedirs(dbDirectory)

rawPrompt = pt.from_template(
    """ 
<s>[INST] You are an expert technical information searching assistant. If the provided information does not have a clear answer, say "I do not have enough information to answer that." [/INST]</s>
[INST] {input}
        Context: {context}
        Answer:
[/INST]
"""
)

class Document:
    def __init__(self, text, metadata=None):
        self.page_content = text
        self.metadata = metadata or {}

# Define functions for each endpoint
def ai_post(query):
    response = llm.invoke(query)
    if "I do not have enough information" in response or "I cannot answer" in response:
        return {"response": "The information provided is insufficient to answer your query accurately."}
    else:
        return {"response": response}

def db_query(query):
    vectorStore = Chroma(persist_directory=dbDirectory, embedding_function=embedding)
    retriever = vectorStore.as_retriever(
        search_type="similarity_score_threshold",
        search_kwargs={"k": 5, "score_threshold": 0.01}
    )
    documentChain = csdc(llm, rawPrompt)
    chain = crc(retriever, documentChain)
    result = chain.invoke({"input": query})

    sources = []
    for doc in result['context']:
        if 'source' in doc.metadata:
            sources.append({
                'source': doc.metadata['source'],
                'Page number': doc.metadata['page'],
                "Page Content": doc.page_content[:50] + "..."
            })

    return {"answer": result['answer'], "sources": sources}

def uploads(files):
    responses = []
    all_chunks = []

    for file in files:
        fileName = file.name.lower()
        file_extension = fileName.split('.')[-1]
        save_directory = None

        if file_extension == 'pdf':
            save_directory = r".\CorporateGPT2\pdfUploads"
        elif file_extension in ['doc', 'docx']:
            save_directory = r".\CorporateGPT2\wordUploads"
        elif file_extension in ['xls', 'xlsx']:
            save_directory = r".\CorporateGPT2\excelUploads"
        elif file_extension in ['ppt', 'pptx']:
            save_directory = r".\CorporateGPT2\powerpointUploads"
        else:
            responses.append({
                "Status": "Unsupported file type",
                "File name": fileName
            })
            continue

        if not os.path.exists(save_directory):
            os.makedirs(save_directory)

        saveFile = os.path.join(save_directory, fileName)
        try:
            with open(saveFile, "wb") as f:
                f.write(file.getbuffer())
            docs = []
            num_pages_or_sheets = 0
            if file_extension == 'pdf':
                loader = pdfPL(saveFile)
                docs = loader.load_and_split()
                num_pages_or_sheets = len(docs)
            elif file_extension in ['doc', 'docx']:
                doc = DocxDocument(saveFile)
                texts = [p.text for p in doc.paragraphs if p.text.strip()]
                docs = [Document(text) for text in texts]
                num_pages_or_sheets = len(docs)
            elif file_extension in ['xls', 'xlsx']:
                excel = pd.ExcelFile(saveFile)
                for sheet_name in excel.sheet_names:
                    sheet_df = pd.read_excel(excel, sheet_name=sheet_name)
                    sheet_text = sheet_df.to_string()
                    docs.append(Document(sheet_text))
                num_pages_or_sheets = len(excel.sheet_names)
            elif file_extension in ['ppt', 'pptx']:
                presentation = Presentation(saveFile)
                for slide in presentation.slides:
                    slide_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                    if slide_text.strip():
                        docs.append(Document(slide_text))
                num_pages_or_sheets = len(presentation.slides)

            chunks = text_splitter.split_documents(docs)
            all_chunks.extend(chunks)

            response = {
                "Status": "Successfully uploaded",
                "File name": fileName,
                "Number of pages or sheets": num_pages_or_sheets,
                "Number of chunks": len(chunks)
            }
            responses.append(response)

        except Exception as e:
            responses.append({
                "Status": "Failed to save file",
                "File name": fileName,
                "Error": str(e)
            })

    return all_chunks, responses

def save_to_db(files):
    chunks, response = uploads(files)
    vectorStore = Chroma.from_documents(documents=chunks, embedding=embedding, persist_directory=dbDirectory)
    vectorStore.persist()
    # Add database upload status to the response
    for i in range(0, len(response)):
        response[i]['Status'] = "Successfully uploaded into Vector database"
    return response

# Load environment variables from .env file
load_dotenv()

# Read database configuration from environment variables
db_file = os.getenv('DB_FILE', 'user_database.db')

# Initialize UserAuth
auth = UserAuth(db_file=db_file)

# Automatically get location data
def get_location_data():
    try:
        ip_response = requests.get("https://api.ipify.org?format=json")
        ip = ip_response.json()["ip"]

        geo_response = requests.get(f"https://ipinfo.io/{ip}/json")
        geo_data = geo_response.json()

        return ip, geo_data
    except Exception as e:
        st.error("Failed to retrieve location data.")
        return None, None

# Function to save location data to the database
def save_location_data(user_id, status, geo_data, ip):
    conn = sqlite3.connect(db_file)
    cursor = conn.cursor()
    
    coordinates = geo_data.get('loc') if geo_data else None
    event_time = datetime.now().isoformat()  # Use ISO format for datetime
    
    cursor.execute("""
        INSERT INTO UserLocationData (user_id, status, event_time, geolocation_coordinates, ip_address)
        VALUES (?, ?, ?, ?, ?)
    """, (user_id, status, event_time, coordinates, ip))
    
    conn.commit()
    cursor.close()
    conn.close()

# Main application logic
if 'authenticated' not in st.session_state:
    st.session_state.authenticated = False
if 'doc_info' not in st.session_state:
    st.session_state.doc_info = None
if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []
if 'user_input' not in st.session_state:
    st.session_state.user_input = ''

if st.session_state.authenticated:
    st.sidebar.write(f"Welcome, {st.session_state.email}")
    
    if 'login_event_saved' not in st.session_state:
        ip, geo_data = get_location_data()
        save_location_data(st.session_state.user_id, "Logged In", geo_data, ip)
        st.session_state.login_event_saved = True
    
    if st.sidebar.button("Logout"):
        ip, geo_data = get_location_data()
        save_location_data(st.session_state.user_id, "Logged Out", geo_data, ip)
        st.session_state.authenticated = False
        st.session_state.email = None
        st.session_state.user_id = None
        st.session_state.reset_code = None
        st.session_state.doc_info = None
        st.session_state.login_event_saved = False
        st.rerun()

    # Sidebar content: source selection, document upload, database, and submit button
    st.sidebar.header("Upload Documents and Connect to Database")
    
    uploaded_files = st.sidebar.file_uploader("Choose a file", type=["docx", "xlsx", "pptx", "pdf"], accept_multiple_files=True)
    database = st.sidebar.text_input("Enter database name")
    submit_button = st.sidebar.button("Submit")

    if submit_button and uploaded_files:
        # Prepare files for upload
        files = []
        for file in uploaded_files:
            files.append(('file', file))

        # Upload the document to the Flask server
        response = save_to_db(files)

        st.sidebar.write(response)

    # Main Page: LLM Chat Interface
    st.title("LLM Chat Interface")

    # Choose between vector database or general model based on user input
    mode = st.selectbox("Choose mode:", ["General Model", "Vector Database"])

    # User input for chat
    user_input = st.text_input("Ask a question:", value=st.session_state.user_input, key="user_input")

    # Buttons to send or cancel the query
    col1, col2 = st.columns([1, 1])
    with col1:
        send_button = st.button("Send", key="send_button")
    with col2:
        cancel_button = st.button("Cancel", key="cancel_button")

    # Display previous chat history with the latest chat on top
    def display_chat_history():
        for chat in reversed(st.session_state.chat_history):
            st.write(f"**User:** {chat['user']}")
            st.write(f"**LLM:** {chat['llm']}")
            st.markdown("---")

    display_chat_history()

    # Clear the input box if the cancel button is clicked
    if cancel_button:
        st.session_state.user_input = ''
        st.rerun()

    # Handle sending the query when the send button is clicked
    if send_button:
        if user_input:
            # Choose between vector database or general model based on user input
            if mode == "Vector Database":            
                answer_data = db_query(user_input)
                answer = answer_data['answer']
                sources = answer_data['sources']

                st.session_state.chat_history.append({"user": user_input, "llm": answer, "sources": sources})

                # Clear the input box after sending
                display_chat_history()
                st.rerun()  # Refresh the app to show the updated chat history

            else:
                response = ai_post(user_input)

                st.session_state.chat_history.append({"user": user_input, "llm": response['response']})

                # Clear the input box after sending
                display_chat_history()
                st.rerun()  # Refresh the app to show the updated chat history

    # Option to clear chat history
    if st.button("Clear Chat History"):
        st.session_state.chat_history = []
        st.rerun()

else:
    st.title("User Authentication")

    auth_option = st.selectbox("Choose option:", ["Log In", "Sign Up"])

    if auth_option == "Log In":
        email = st.text_input("Email")
        password = st.text_input("Password", type="password")

        if st.button("Log In"):
            user_id = auth.login_user(email, password)
            if user_id:
                st.session_state.authenticated = True
                st.session_state.email = email
                st.session_state.user_id = user_id
                st.success("Logged in successfully!")
                st.rerun()
            else:
                st.error("Invalid email or password.")

    elif auth_option == "Sign Up":
        email = st.text_input("Email")
        password = st.text_input("Password", type="password")
        confirm_password = st.text_input("Confirm Password", type="password")
        first_name = st.text_input("Enter First Name")
        surname = st.text_input("Enter Surname")
        mobile_number = st.text_input("Enter Mobile Number")

        if st.button("Sign Up"):
            if password == confirm_password:
                user_id = auth.register_user(email, password, first_name, surname, mobile_number)
                if user_id:
                    st.success('Signed up successfully! Please select "log in" option above to log in.', icon="🟢")
                else:
                    st.error("🔴 Failed to sign up. Email might already be registered.")
            else:
                st.error("🔴 Passwords do not match.")

    elif auth_option == "Reset Password":
        email = st.text_input("Email")

        if st.button("Send Reset Code"):
            reset_code = auth.send_reset_code(email)
            if reset_code:
                st.success("Reset code sent to your email.")
                st.session_state.reset_code = reset_code
                st.session_state.email = email
                st.session_state.authenticated = True
                st.rerun()
            else:
                st.error("Failed to send reset code.")
