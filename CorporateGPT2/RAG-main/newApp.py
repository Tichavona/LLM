from langchain_community.llms import Ollama
from flask import Flask, request, jsonify
from langchain.text_splitter import RecursiveCharacterTextSplitter as rcts
from langchain_community.embeddings.fastembed import FastEmbedEmbeddings as fee
# from langchain_community.embeddings.bedrock import BedrockEmbeddings as be
from langchain_community.embeddings.ollama import OllamaEmbeddings as oe
from langchain_community.document_loaders import PDFPlumberLoader as pdfPL
from langchain_community.vectorstores import Chroma
from langchain.chains.combine_documents import create_stuff_documents_chain as csdc
from langchain.chains import create_retrieval_chain as crc
from langchain.prompts import PromptTemplate as pt
from docx import Document as DocxDocument
from langchain.docstore.document import Document
import pandas as pd
from pptx import Presentation
import os

app = Flask(__name__)

llm = Ollama(model="llama3")

embedding = fee()

# def get_embedding_function():
#     embeddings = be()
#     return embeddings

def get_embedding_function():
    embeddings = oe(model="nomic-embed-text")
    return embeddings

text_splitter = rcts(
    chunk_size=2048,
    chunk_overlap=200,
    length_function=len,
    is_separator_regex='\n'
)

# Ensure the database directory exists or create it
dbDirectory = r"C:\Users\ACER\Documents\DecisionEdgeConcepts\CorporateGPT2\db"
if not os.path.exists(dbDirectory):
    os.makedirs(dbDirectory)

rawPrompt = pt.from_template(
""" 
<s>[INST] You are an expert technical information searching assistant. If the provided information does not have a clear answer, say "I do not have enough information to answer that." [/INST]</s>
[INST] {input}
        Context: {context}
        Answer:
[/INST]
"""
)

class Document:
    def __init__(self, text, metadata=None):
        self.page_content = text
        self.metadata = metadata or {}

@app.route("/ai", methods=['POST'])
def aiPost():
    print('Post /ai called')
    json_content = request.json
    query = json_content.get("query")
    print(f'query: {query}')

    response = llm.invoke(query)
    print(response)

    # Check if the response is relevant
    if "I do not have enough information" in response or "I cannot answer" in response:
        responseAnswer = {"response": "The information provided is insufficient to answer your query accurately."}
    else:
        responseAnswer = {"response": response}
    
    return jsonify(responseAnswer)

@app.route("/dbQuerying", methods=['POST'])
def dbQuery():
    print('Post /Querying documents in Vector database')
    json_content = request.json
    query = json_content.get("query")
    print(f'query: {query}')

    # Loading Vectorstore database
    vectorStore = Chroma(persist_directory=dbDirectory, embedding_function=fee)#embedding_function=embedding)

    # Creating retriever
    retriever = vectorStore.as_retriever(
        search_type="similarity_score_threshold",
        search_kwargs={
            "k": 5,
            "score_threshold": 0.01  # Stricter threshold
        }
    )

    documentChain = csdc(llm, rawPrompt)

    chain = crc(retriever, documentChain)

    result = chain.invoke({"input": query})

    print(result)

    sources = []
    for doc in result['context']:
        if 'source' in doc.metadata:
            sources.append(
                {
                    'source' : doc.metadata['source'],
                    'Page number' : doc.metadata['page'],
                    "Page Content" : doc.page_content[:50] + "..."
                }
            )

    response = {"answer": result['answer'], "sources" : sources}
    return jsonify(response)

@app.route("/Uploads", methods=['POST'])
def Uploads():
    if 'file' not in request.files:
        print(request.files)
        return None, {"Status": "No files found in the request"}, 400
    
    files = request.files.getlist('file')

    if not files or all(file.filename == '' for file in files):
        return None, {"Status": "No files selected"}, 400

    responses = []
    all_chunks = []

    for file in files:
        fileName = file.filename.lower()
        file_extension = fileName.split('.')[-1]

        # Directory to save the uploaded file
        save_directory = None
        if file_extension == 'pdf':
            save_directory = r"C:\Users\ACER\Documents\DecisionEdgeConcepts\CorporateGPT2\pdfUploads"
        elif file_extension in ['doc', 'docx']:
            save_directory = r"C:\Users\ACER\Documents\DecisionEdgeConcepts\CorporateGPT2\wordUploads"
        elif file_extension in ['xls', 'xlsx']:
            save_directory = r"C:\Users\ACER\Documents\DecisionEdgeConcepts\CorporateGPT2\excelUploads"
        elif file_extension in ['ppt', 'pptx']:
            save_directory = r"C:\Users\ACER\Documents\DecisionEdgeConcepts\CorporateGPT2\powerpointUploads"
        else:
            responses.append({
                "Status": "Unsupported file type",
                "File name": fileName
            })
            continue

        if not os.path.exists(save_directory):
            os.makedirs(save_directory)

        saveFile = os.path.join(save_directory, fileName)

        try:
            # Save the file
            file.save(saveFile)
            print(f'File name: {fileName}')

            # Load file based on type
            docs = []
            num_pages_or_sheets = 0
            if file_extension == 'pdf':
                loader = pdfPL(saveFile)
                docs = loader.load_and_split()
                num_pages_or_sheets = len(docs)
            elif file_extension in ['doc', 'docx']:
                doc = DocxDocument(saveFile)
                texts = [p.text for p in doc.paragraphs if p.text.strip()]
                docs = [Document(text) for text in texts]
                num_pages_or_sheets = len(docs)
            elif file_extension in ['xls', 'xlsx']:
                excel = pd.ExcelFile(saveFile)
                for sheet_name in excel.sheet_names:
                    sheet_df = pd.read_excel(excel, sheet_name=sheet_name)
                    sheet_text = sheet_df.to_string()
                    docs.append(Document(sheet_text))
                num_pages_or_sheets = len(excel.sheet_names)
            elif file_extension in ['ppt', 'pptx']:
                presentation = Presentation(saveFile)
                for slide in presentation.slides:
                    slide_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                    if slide_text.strip():
                        docs.append(Document(slide_text))
                num_pages_or_sheets = len(presentation.slides)

            chunks = text_splitter.split_documents(docs)
            all_chunks.extend(chunks)

            # User feedback of process completion status if successful
            response = {
                "Status": "Successfully uploaded",
                "File name": fileName,
                "Number of pages or sheets": num_pages_or_sheets,
                "Number of chunks": len(chunks)
            }
            responses.append(response)

        except Exception as e:
            responses.append({
                "Status": "Failed to save file",
                "File name": fileName,
                "Error": str(e)
            })

    if not all_chunks:
        return None, {"Status": "No files processed successfully"}, 500

    return all_chunks, responses, 200

def add_to_chroma(new_chunks, new_chunk_ids):
    db = Chroma(
        persist_directory=dbDirectory,
        embedding_function=fee
    )
    print(new_chunks, new_chunk_ids)
    db.add_documents(documents = new_chunks, id=new_chunk_ids)
    
@app.route("/saveDB", methods=['POST'])
def save_to_db():
    chunks, response, status_code = Uploads()

    if status_code != 200:
        return jsonify(response), status_code
    
    last_page_id = None
    current_chunk_index = 0

    chunksFile = r"C:\Users\ACER\Documents\DecisionEdgeConcepts\CorporateGPT2\RAG-main\allChunks.csv"

    if os.path.exists(chunksFile):
        allChunks = pd.read_csv(chunksFile)
        all_chunk_ids = list(allChunks["myChunks"])
    else:
        all_chunk_ids = []

    new_chunks = []
    new_chunk_ids = []

    for chunk in chunks:
        source = chunk.metadata.get("source")
        page = chunk.metadata.get("page")
        current_page_id = f"{source}:{page}"
        if current_page_id == last_page_id:
            current_chunk_index += 1
        else:
            current_chunk_index = 0
        chunk_id = f"{current_page_id}:{current_chunk_index}"
        chunk.metadata["id"] = chunk_id
        last_page_id = current_page_id
        if chunk_id not in all_chunk_ids:
            new_chunks.append(chunk)
            new_chunk_ids.append(chunk_id)
            all_chunk_ids.append(chunk_id)        

    # Save chunks to the vector database
    for i in range(0, len(new_chunks)):
        vectorStore = add_to_chroma(new_chunks[i], new_chunk_ids[i])
    vectorStore.persist()

    allChunks = pd.DataFrame(all_chunk_ids, columns=["myChunks"])
    allChunks.to_csv(chunksFile, index=False)

    # Add database upload status to the response
    for i in range(0, len(response)):
        response[i]['Status'] = "Successfully uploaded into Vector database"

    # Return the updated response
    return jsonify(response), 200

@app.route("/", methods = ['GET', 'POST'])
def home():
    response = """
    Welcome to the RAG API. Use the specific endpoints to interact with the system. The API has the following functions
    1. Post to Vector database
    2. Chat with general model
    3. Chat with the corporate data
    """
    return response

def startApp():
    app.run(host="0.0.0.0", port=8080, debug=True)

if __name__ == "__main__":
    startApp()
